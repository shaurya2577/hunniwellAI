#!/usr/bin/env python3
"""
ingest.py — Company Files → Airtable ingestion tool

Walks $HUNNIWELL_COMPANYFILES_ROOT/<EVENT>/<COMPANY>/, reads every PDF / DOCX /
PPTX / TXT / MD inside each company folder, asks Claude to extract a strict-
schema JSON record, and POSTs it to Airtable. Re-runs skip companies already
recorded in .processed.json. EVENT_LAYOUTS lives in events.py (edit there to
add new events).

Setup:
    cd /Users/sbhartia/Dev/Hunniwell
    python3 -m venv .venv && source .venv/bin/activate
    pip install -e .                      # installs all repo deps
    cp ai/airtable_ingest/.env.example ai/airtable_ingest/.env  # then fill keys

Required env vars (loaded from .env next to ingest.py):
    ANTHROPIC_API_KEY
    AIRTABLE_API_KEY                Airtable personal access token
    AIRTABLE_BASE_ID                e.g. app...
    AIRTABLE_TABLE_NAME             table name or table ID (tbl...)
    HUNNIWELL_COMPANYFILES_ROOT     path to CompanyFiles tree (default: ~/Documents/Hunniwell)

Run (from repo root):
    python -m ai.airtable_ingest.ingest                       # full run
    python -m ai.airtable_ingest.ingest --dry-run             # no Airtable writes
    python -m ai.airtable_ingest.ingest --event "JPM 2026 (260115)"
    python -m ai.airtable_ingest.ingest --company "Auvi Labs"
    python -m ai.airtable_ingest.ingest --root /custom/path
    python -m ai.airtable_ingest.ingest --reset-state
"""

import argparse
import csv
import json
import os
import re
import sys
import time
from pathlib import Path
from urllib.parse import quote

import anthropic
import requests
from docx import Document
from dotenv import load_dotenv
from pptx import Presentation
from pypdf import PdfReader

SCRIPT_DIR = Path(__file__).resolve().parent
load_dotenv(SCRIPT_DIR / ".env")

# Default location of the CompanyFiles tree on this user's machine. Override at
# runtime via --root PATH or the HUNNIWELL_COMPANYFILES_ROOT env var.
DEFAULT_ROOT = Path("/Users/sbhartia/Documents/Hunniwell")
STATE_FILE = SCRIPT_DIR / ".processed.json"
LOG_FILE = SCRIPT_DIR / "run_log.csv"
ERRORS_DIR = SCRIPT_DIR / "errors"

CLAUDE_MODEL = "claude-haiku-4-5-20251001"  # was claude-sonnet-4-5; Haiku is ~5x faster, ~10x cheaper, fine for this domain
DATA_ENTRY = "AI"

PER_FILE_CHAR_LIMIT = 30_000
TOTAL_CHAR_LIMIT = 120_000

READABLE_EXTS = {".pdf", ".docx", ".pptx", ".txt", ".md"}

# JSON key from Claude -> exact Airtable field name
FIELD_MAP = {
    "company": "Company",
    "event": "Event",
    "data_entry": "Data Entry",
    "iata_code": "IATA Code",
    "country": "Country",
    "short_description": "Short Description / Key technology",
    "medical_field": "Medical Field",
    "indication": "Indication",
    "class_of_device": "Class of Device",
    "regulatory_pathway": "Regulatory Pathway",
    "dev_stage": "Dev. Stage",
    "dev_stage_details": "Dev Stage Details",
    "long_description": "Long Description (medical/clinical)",
    "equity_raised_m": "Equity Raised ($M)",
    "coming_round": "Coming Round",
    "size_of_round_m": "Size of Round ($M)",
    "est_close": "Est. Close (Month/Year)",  # note: Airtable field has a trailing space
    "key_executives": "Key Executive(s)",
    "ceo_email": "CEO's Email",
    "ceo_cell": "CEO's Cell #",
    "company_notes": "Company Notes",
    "url": "Url",
    "address": "Address",
}

# EVENT_LAYOUTS lives in events.py so non-coders can add new events without touching ingest.py.
# Folders not listed there are skipped when --event is not given.
from .events import EVENT_LAYOUTS  # noqa: E402

SYSTEM_PROMPT = """You are extracting structured data about a single medtech company from raw text dumped out of the company's deck and any summary documents. Return only a JSON object matching the schema below — no prose, no markdown fences.

Strict rule: If the provided text does not clearly support a field, omit that key entirely. Do not infer, do not guess, do not fill from general knowledge. Cite sources where possible. Leaving a field blank is correct and expected. Do not repeat information across fields; each field should capture unique data points. 
Focus on precision and verifiability based on the text, not completeness. Remove connecting words; uses phrases to make fields as concise as possible while still being clear.

CRITICAL CITATION RULES:
1. The source files contain explicit inline source brackets (e.g., [INT-1], [INT-2], [EXT-4], [EXT-11]).
2. For text fields—specifically 'short_description' and 'company_notes'—you MUST append the exact, corresponding source bracket to the end of the sentence or metric you extract.
3. Example format: "The company completed its first GMP production run in Q2 2022 [INT-1]. They are targeting a $4.0M Series A round [INT-1]."
4. Never strip away or omit these source brackets; they are critical for auditing.

Schema keys (all values must be strings):
iata_code, country, short_description, medical_field, indication, class_of_device, regulatory_pathway, dev_stage, dev_stage_details, long_description, equity_raised_m, coming_round, size_of_round_m, est_close, key_executives, ceo_email, ceo_cell, company_notes, url, address.

- address: Refer to online sources if needed. Full address of company headquarters, ideally including street, city, state, and zip (if US-based).
- iata code must be of closest international airport to company hq address
- short_description must be 2-4 sentences and MUST STRICTLY follow the structure of:
"[Technology name] provides a way to solve [Problem] in [Population] in order to [Outcome]." We want to focus on the companies/technology's target population and outcome underpinning the 
technology and what makes it better than other technologies out there . There is no need to repeat the company's name in the short description. Here are some examples of what we want a short description looks like: [Technology name] is a way to 
relieve urinary symptoms in men with AUASI score > 14 BPH-related urinary retention that has greater effectiveness than current minimally invasive treatments and has fewer complications than surgical treatment options, [Technology name] is a way
 to perform point-of-care testing in patients  with skin lesions in order to enable accurate, inexpensive diagnosis of malignant melanoma by a dermatologist.

- medical_field and indication: comma-separated string if multiple values apply (e.g. "Nephrology, Vascular Surgery, Dialysis").
- equity_raised_m and size_of_round_m: decimal millions as a string (e.g. "0.10", "1.0").
- est_close: format MM/YYYY.
- company_notes: Capture high-signal, diligence-relevant metrics not covered elsewhere in this profile. Prioritize: Key Institutional 
Investors/Lead Backers, IP Portfolio Strength, Prior Marketed Products/Track Record, Current Headcount.Strict Constraints: Focus on financial strategic, or technical moats. Do not include employee count. Where precise internal metrics are unavailable,
 reference or anchor the data to verified online sources, industry databases, or public filings. Do not include team size. Leave this field blank or omit it completely if no unique, 
 high-signal data exists beyond what has already been captured. Structure: a) awards/prizes/achievements for the company or the founders
b) non-dilutive grants c) rounds of investments - $ amt, when, valuation & identity of investors where available
d) Direct observations and comments/feedback from reliable sources
- long_description: The Long Description must be 4-6 sentences and provide comprehensive clinical/medical context. 
    STRUCTURE: Sentence 1: "[Technology name] is a [category: device/diagnostic/software/therapeutic] that addresses [clinical problem] in [target population] by [mechanism of action/how it works]." Sentence 2-3: Clinical differentiation, explain what makes this better:
     Clinical evidence: "Published data shows [clinical outcome metric] (e.g., 2.5x faster healing, 40% reduction in complications)"Competitive advantage: "Unlike [competitor/current standard of care], [technology name] [specific advantage]." 
     Market fit: "[Population size] patients in [geography] currently rely on [existing solution], creating a [market size] opportunity." Sentences 4-5: Regulatory and commercial status: 
     Current stage: "The technology has [regulatory clearance: FDA 510(k)/De Novo/PMA/CE Mark/CLIA/etc.] and is [status: in development/clinical trials/commercially available]." Revenue/traction: "[Number] procedures/patients treated to date" OR "Early access/pilot programs in [location]." 
    RULES: Do NOT repeat the company name except in first sentence. Do NOT infer clinical outcomes; only include published/presented data from the documents.
    Do NOT guess regulatory status; only include what is explicitly stated. Omit any sentence component if data is not clearly present in documents. Focus on what DIFFERENTIATES this technology, not generic features. If clinical evidence is limited, emphasize regulatory progress or market opportunity instead.
    Use specific numbers (percentages, patient counts, pricing) when available. Reference "current standard of care" or specific competitor names when available 

"""

CHALLENGER_PROMPT = """You are a strict, adversarial Venture Capital Auditor and Copyeditor. You are given raw text documents from a medtech company's pitch materials and an initial AI data entry draft.

Your two-part task is to:
1. AUDIT: Eliminate hallucinations and unsupported claims.
2. EDIT FOR READABILITY: Ensure the text is clean, professional, and punchy—preserving all exact metrics and citation brackets while stripping away dense academic jargon or clunky phrasing.

<editorial_rules>
- COMPACTNESS: Use active, concise phrases instead of long connecting clauses. Do not repeat the company's name unless strictly required by a specific schema layout. 
- SHORT_DESCRIPTION TEMPLATE: Force the 'short_description' field to read as an investor-facing value proposition. Ensure it highlights what makes it better than the standard of care. Be concise.
- MEDICAL_FIELD CLAUSE: Ensure 'medical_field' preserves high-level industry taxonomies (e.g., if it narrows down to "In Vitro Diagnostics", make sure broader context tags like "Oncology" or "Gastroenterology" remain intact so database filtering works).
- READABILITY OVER CLUTTER: Do not let definitions spill into multi-paragraph run-ons. If a text field contains multiple distinct metrics, combine them elegantly into 2-3 highly polished sentences.
- PRESERVE ALL CITATIONS: You MUST preserve all bracketed citations (e.g., [INT-1], [EXT-4]) exactly where their corresponding factual statement or metric lands. Never strip them out.
</editorial_rules>

<audit_rules>
- If any info in the initial draft is NOT explicitly supported or cited by the raw text, wipe that value completely (set to "").
- If a financial or regulatory number doesn't match the raw documents perfectly, correct it to match the truth.
- Be as concise as possible while still conveying all the key data points. Remove any fluff or filler words that don't add concrete information.
</audit_rules>

STRICT FORMATTING CONSTRAINT: 
You must output ONLY a valid JSON object matching the exact schema keys provided in the draft. Do NOT include any introductory text, closing notes, markdown blocks, or explanations. Start your response directly with '{' and end it directly with '}'.

"""

def derive_event(folder_name: str) -> str:
    """Canonical folder names ARE the Airtable Event values — no transformation."""
    return folder_name.strip()


def extract_pdf(path: Path) -> str:
    try:
        reader = PdfReader(str(path))
        parts = []
        for page in reader.pages:
            try:
                parts.append(page.extract_text() or "")
            except Exception:
                continue
        return "\n".join(parts)
    except Exception as e:
        print(f"  [warn] PDF extract failed for {path.name}: {e}")
        return ""


def extract_docx(path: Path) -> str:
    try:
        doc = Document(str(path))
        parts = [p.text for p in doc.paragraphs if p.text]
        for tbl in doc.tables:
            for row in tbl.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n".join(parts)
    except Exception as e:
        print(f"  [warn] DOCX extract failed for {path.name}: {e}")
        return ""


def extract_pptx(path: Path) -> str:
    try:
        prs = Presentation(str(path))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"--- Slide {i} ---")
            for shape in slide.shapes:
                try:
                    if shape.has_text_frame:
                        text = shape.text_frame.text
                        if text:
                            parts.append(text)
                except Exception:
                    continue
            try:
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text
                    if notes:
                        parts.append(f"[notes] {notes}")
            except Exception:
                pass
        return "\n".join(parts)
    except Exception as e:
        print(f"  [warn] PPTX extract failed for {path.name}: {e}")
        return ""


def extract_text_file(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8", errors="ignore")
    except Exception as e:
        print(f"  [warn] text read failed for {path.name}: {e}")
        return ""


EXTRACTORS = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".txt": extract_text_file,
    ".md": extract_text_file,
}


def collect_company_text(company_dir: Path) -> str:
    files: list[Path] = []
    for path in company_dir.rglob("*"):
        if not path.is_file():
            continue
        rel_depth = len(path.relative_to(company_dir).parts)
        if rel_depth > 2:
            continue
        if path.suffix.lower() in READABLE_EXTS:
            files.append(path)

    chunks: list[str] = []
    total = 0
    for f in sorted(files):
        text = EXTRACTORS[f.suffix.lower()](f)
        if not text.strip():
            continue
        text = text[:PER_FILE_CHAR_LIMIT]
        header = f"\n\n=== {f.relative_to(company_dir)} ===\n"
        remaining = TOTAL_CHAR_LIMIT - total
        if remaining <= len(header):
            break
        if len(header) + len(text) > remaining:
            text = text[: remaining - len(header)]
        chunks.append(header + text)
        total += len(header) + len(text)
        if total >= TOTAL_CHAR_LIMIT:
            break
    return "".join(chunks)


def _strip_fences(s: str) -> str:
    s = s.strip()
    if s.startswith("```"):
        s = re.sub(r"^```(?:json)?\s*", "", s)
        s = re.sub(r"\s*```$", "", s)
    return s.strip()


def _coerce_value(v) -> str:
    if v is None:
        return ""
    if isinstance(v, list):
        return ", ".join(str(x).strip() for x in v if str(x).strip())
    if isinstance(v, (dict,)):
        return json.dumps(v, ensure_ascii=False)
    return str(v).strip()


def call_claude(client: anthropic.Anthropic, company_name: str, text_blob: str) -> dict:
    user_msg = f"Company folder name: {company_name}\n\n{text_blob}"
    last_transport_err = None
    raw_content = ""
    for attempt in range(2):
        try:
            resp = client.messages.create(
                model=CLAUDE_MODEL,
                max_tokens=2000,
                system=SYSTEM_PROMPT,
                messages=[{"role": "user", "content": user_msg}],
                timeout=120.0,
            )
            text_pieces = [
                getattr(b, "text", "")
                for b in getattr(resp, "content", [])
                if getattr(b, "type", None) == "text" and getattr(b, "text", None)
            ]
            if not text_pieces:
                stop_reason = getattr(resp, "stop_reason", None)
                raise ValueError(f"Claude returned no text content (stop_reason={stop_reason})")
            raw_content = "".join(text_pieces)
            cleaned = _strip_fences(raw_content)
            parsed = json.loads(cleaned)
            if not isinstance(parsed, dict):
                raise ValueError(
                    f"Claude returned non-object JSON ({type(parsed).__name__}); "
                    f"first 300 chars: {str(parsed)[:300]}"
                )
            return parsed
        except (anthropic.APIConnectionError, anthropic.APITimeoutError, anthropic.RateLimitError) as e:
            last_transport_err = e
            if attempt == 0:
                time.sleep(5)
                continue
            raise
        except json.JSONDecodeError as e:
            raise ValueError(f"Claude returned non-JSON. Raw: {raw_content[:1000]}") from e
    if last_transport_err:
        raise last_transport_err
    raise RuntimeError("call_claude: unreachable")

def call_challenger(client: anthropic.Anthropic, text_blob: str, draft_record: dict) -> dict:
    user_msg = f"""RAW COMPANY DOCUMENTS:\n{text_blob}\n\n
INITIAL AI DRAFT:\n{json.dumps(draft_record, indent=2)}"""
    
    try:
        resp = client.messages.create(
            model="claude-sonnet-4-6", 
            max_tokens=2500, 
            system=CHALLENGER_PROMPT,
            messages=[{"role": "user", "content": user_msg}],
            timeout=150.0, 
        )
        raw_content = "".join([b.text for b in resp.content if b.type == "text"]) 
       
        start_idx = raw_content.find('{')
        end_idx = raw_content.rfind('}')
        
        if start_idx != -1 and end_idx != -1 and end_idx > start_idx:
            cleaned = raw_content[start_idx:end_idx + 1]
        else:
            cleaned = _strip_fences(raw_content) 
            
        return json.loads(cleaned) 
    except Exception as e:
        print(f"  [warn] Adversarial challenge failed, falling back to original draft: {e}") 
        return {} 


def write_airtable(record: dict, base_id: str, table: str, api_key: str) -> str:
    fields = {}
    for json_key, val in record.items():
        airtable_name = FIELD_MAP.get(json_key)
        if not airtable_name:
            continue
        coerced = _coerce_value(val)
        if not coerced:
            continue
        fields[airtable_name] = coerced
    payload = { "records": [{"fields": fields}],
    "typecast": True, 
    "performUpsert": {"fieldsToMergeOn": ["Company","Event"]}
    }

    url = f"https://api.airtable.com/v0/{base_id}/{quote(table, safe='')}"
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
    }

    for attempt in range(2):
        r = requests.patch(url, headers=headers, json=payload, timeout=30)
        if r.status_code in (200, 201):
            # return r.json()["id"]
            return r.json()["records"][0]["id"]
        if r.status_code in (429, 500, 502, 503, 504) and attempt == 0:
            time.sleep(5)
            continue
        raise RuntimeError(f"Airtable {r.status_code}: {r.text}")
    raise RuntimeError("Airtable retries exhausted")


def load_state(path: Path) -> set[str]:
    if path.exists():
        try:
            return set(json.loads(path.read_text()))
        except Exception:
            return set()
    return set()


def save_state(state: set[str], path: Path) -> None:
    path.write_text(json.dumps(sorted(state), indent=2))


def iter_event_companies(event_folder: Path):
    """Yield (company_dir, state_key, event_label) for an event folder.

    Honors EVENT_LAYOUTS entries (`{"layout": "flat"|"nested", "categories": [...]}`).
    Folders absent from EVENT_LAYOUTS default to flat with no category validation."""
    entry = EVENT_LAYOUTS.get(event_folder.name, {"layout": "flat"})
    layout = entry.get("layout", "flat")
    declared_categories = entry.get("categories")
    base_event = derive_event(event_folder.name)

    if layout == "nested":
        on_disk_cats = sorted(p for p in event_folder.iterdir() if p.is_dir())
        if declared_categories:
            declared = set(declared_categories)
            on_disk = {c.name for c in on_disk_cats}
            extras = on_disk - declared
            if extras:
                print(f"  WARN: {event_folder.name}: unexpected category subdirs not in events.py: {sorted(extras)}")
        for cat in on_disk_cats:
            for comp in sorted(p for p in cat.iterdir() if p.is_dir()):
                yield (
                    comp,
                    f"{event_folder.name}/{cat.name}/{comp.name}",
                    f"{base_event} - {cat.name}",
                )
    else:  # flat
        for comp in sorted(p for p in event_folder.iterdir() if p.is_dir()):
            yield (
                comp,
                f"{event_folder.name}/{comp.name}",
                base_event,
            )


def main() -> None:
    ap = argparse.ArgumentParser(description="Walk CompanyFiles/ and ingest each company into Airtable.")
    ap.add_argument("--dry-run", action="store_true", help="Print JSON instead of POSTing to Airtable.")
    ap.add_argument("--event", help="Restrict to a single event folder (exact folder name). Bypasses the EVENT_LAYOUTS whitelist.")
    ap.add_argument("--company", help="Restrict to a single company folder (exact folder name).")
    ap.add_argument("--reset-state", action="store_true", help="Delete the state file and exit.")
    ap.add_argument("--state-file", help="Override state file path (default: .processed.json next to ingest.py). Use a separate path to run parallel processes safely.")
    ap.add_argument("--log-file", help="Override run-log CSV path (default: run_log.csv next to ingest.py).")
    ap.add_argument("--root", help="Override CompanyFiles root (default: $HUNNIWELL_COMPANYFILES_ROOT, then /Users/piusamartey/hunniwellAI).")
    args = ap.parse_args()

    state_path = Path(args.state_file).resolve() if args.state_file else STATE_FILE
    log_path = Path(args.log_file).resolve() if args.log_file else LOG_FILE
    root = Path(args.root or os.environ.get("HUNNIWELL_COMPANYFILES_ROOT") or DEFAULT_ROOT).expanduser()

    if args.reset_state:
        if state_path.exists():
            state_path.unlink()
            print(f"Removed {state_path}")
        else:
            print("No state file to remove.")
        return

    if not root.exists() or not root.is_dir():
        print(f"ERROR: {root} does not exist or is not a directory.")
        sys.exit(1)

    anthropic_key = os.environ.get("ANTHROPIC_API_KEY")
    airtable_key = os.environ.get("AIRTABLE_API_KEY")
    base_id = os.environ.get("AIRTABLE_BASE_ID")
    table_name = os.environ.get("AIRTABLE_TABLE_NAME")

    if not anthropic_key:
        print("ERROR: ANTHROPIC_API_KEY not set.")
        sys.exit(1)
    if not args.dry_run and not all([airtable_key, base_id, table_name]):
        print("ERROR: AIRTABLE_API_KEY / AIRTABLE_BASE_ID / AIRTABLE_TABLE_NAME must be set for non-dry runs.")
        sys.exit(1)

    client = anthropic.Anthropic(api_key=anthropic_key)
    state = load_state(state_path)
    ERRORS_DIR.mkdir(exist_ok=True)

    log_exists = log_path.exists()
    log_fh = log_path.open("a", newline="", encoding="utf-8")
    log_writer = csv.writer(log_fh)
    if not log_exists:
        log_writer.writerow(
            ["company_folder", "event", "status", "airtable_record_id", "error", "fields_populated", "fields_left_blank"]
        )

    all_event_folders = sorted(p for p in root.iterdir() if p.is_dir())
    if args.event:
        event_folders = [p for p in all_event_folders if p.name == args.event]
        if not event_folders:
            print(f"No event folder matching '{args.event}'.")
            log_fh.close()
            sys.exit(1)
    else:
        # Default: only process folders explicitly whitelisted in EVENT_LAYOUTS.
        event_folders = [p for p in all_event_folders if p.name in EVENT_LAYOUTS]
        skipped = [p.name for p in all_event_folders if p.name not in EVENT_LAYOUTS]
        if skipped:
            print(f"Skipping {len(skipped)} non-whitelisted top-level folders: {', '.join(skipped[:8])}{'…' if len(skipped) > 8 else ''}")

    total_ok = 0
    total_skip = 0
    total_fail = 0

    for event_folder in event_folders:
        for company_dir, state_key, event_name in iter_event_companies(event_folder):
            if args.company and company_dir.name != args.company:
                continue
            company = company_dir.name.strip()
            print(f"\n[{event_name}] {company}")

            #if state_key in state:
            #    print("  skip (already processed)")
            #    total_skip += 1
            #    continue

            text_blob = collect_company_text(company_dir)
            if not text_blob.strip():
                print("  skip (no readable content)")
                log_writer.writerow([state_key, event_name, "no_content", "", "", "", ""])
                log_fh.flush()
                total_skip += 1
                continue

            try:
                record = call_claude(client, company, text_blob)

                print("  → Submitting full entry to Adversarial Council for audit...")
                audited_record = call_challenger(client, text_blob, record)
                
                if audited_record and isinstance(audited_record, dict):
                    audited_record["event"] = event_name
                    audited_record["data_entry"] = DATA_ENTRY
                    audited_record["company"] = company
                    record = audited_record
                    print("  → Audit complete. Non-supported fields cleared.")

            except Exception as e:
                err = str(e)
                print(f"  ERROR (claude): {err[:200]}")
                safe = re.sub(r"[^A-Za-z0-9._-]+", "_", company)[:80] or "company"
                (ERRORS_DIR / f"{safe}.txt").write_text(err)
                log_writer.writerow([state_key, event_name, "claude_error", "", err[:500], "", ""])
                log_fh.flush()
                total_fail += 1
                continue

            record["event"] = event_name
            record["data_entry"] = DATA_ENTRY
            record["company"] = company

            populated_keys = [k for k in FIELD_MAP if _coerce_value(record.get(k))]
            blank_keys = [k for k in FIELD_MAP if k not in populated_keys]

            if args.dry_run:
                preview = {FIELD_MAP[k]: _coerce_value(record[k]) for k in populated_keys}
                print("  [dry-run] would post:")
                for line in json.dumps(preview, indent=2, ensure_ascii=False).splitlines():
                    print(f"  {line}")
                log_writer.writerow(
                    [state_key, event_name, "dry_run", "", "", ";".join(populated_keys), ";".join(blank_keys)]
                )
                log_fh.flush()
                continue

            try:
                rec_id = write_airtable(record, base_id, table_name, airtable_key)
                print(f"  -> Airtable {rec_id}   populated={len(populated_keys)} blank={len(blank_keys)}")
                state.add(state_key)
                save_state(state, state_path)
                log_writer.writerow(
                    [state_key, event_name, "ok", rec_id, "", ";".join(populated_keys), ";".join(blank_keys)]
                )
                total_ok += 1
            except Exception as e:
                err = str(e)
                print(f"  ERROR (airtable): {err[:200]}")
                log_writer.writerow(
                    [state_key, event_name, "airtable_error", "", err[:500], ";".join(populated_keys), ";".join(blank_keys)]
                )
                total_fail += 1
            log_fh.flush()

    log_fh.close()
    print(f"\nDone. ok={total_ok} skipped={total_skip} failed={total_fail}")


if __name__ == "__main__":
    main()
