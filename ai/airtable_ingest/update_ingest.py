#!/usr/bin/env python3
"""
update_ingest.py — updates web_injest quarterly 

For each company folder under $HUNNIWELL_COMPANYFILES_ROOT, this script:
  1. Reads all internal docs (PDF / DOCX / PPTX / TXT / MD) as context
  2. Uses Claude Haiku's built-in web_search tool to find additional info online
  3. Writes a cited .txt research report into each company folder

Citation format:
  [INT-1], [INT-2]... = internal files already in the company folder
  [EXT-1], [EXT-2]... = external web sources found by Claude

Output location (default):
  $HUNNIWELL_COMPANYFILES_ROOT/<EVENT>/<COMPANY>/web_research.txt
  Pass --output-dir to collect all reports in one folder instead.

Setup (same venv as ingest.py):
    pip install anthropic python-dotenv pypdf python-docx python-pptx

Required env vars (same .env as ingest.py):
    ANTHROPIC_API_KEY
    HUNNIWELL_COMPANYFILES_ROOT

Run:
    python web_ingest.py                                # all companies (skips up-to-date)
    python web_ingest.py --event "JPM 2026 (260115)"
    python web_ingest.py --company "Auvi Labs"
    python web_ingest.py --company "Auvi Labs" --force  # re-run even if up-to-date
    python web_ingest.py --dry-run                      # print reports, don't write files
    python web_ingest.py --reset-state                  # clear skip cache
    python web_ingest.py --output-dir ./reports         # write all .txt to one folder
    python web_ingest.py --stale-months 3               # override staleness threshold (default: 3)

Update logic (runs every quarter automatically):
  A company is re-researched if ANY of the following are true:
    1. Its web_ingest.txt is older than --stale-months (default 3)
    2. Any internal doc in the company folder is newer than the existing report
    3. --force flag is passed (always re-runs the targeted company/event)
  If none of the above: skip (already up-to-date).
"""

import argparse
import json
import os
import re
import sys
import time
from datetime import datetime, timezone
from pathlib import Path

import anthropic
from dotenv import load_dotenv
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

# ---------------------------------------------------------------------------
# Config
# ---------------------------------------------------------------------------

SCRIPT_DIR = Path(__file__).resolve().parent
load_dotenv(SCRIPT_DIR / ".env")

DEFAULT_ROOT = Path(os.environ.get("HUNNIWELL_COMPANYFILES_ROOT", "~/Documents/Hunniwell")).expanduser()
STATE_FILE = SCRIPT_DIR / ".web_processed.json"

CLAUDE_MODEL = "claude-haiku-4-5-20251001"

PER_FILE_CHAR_LIMIT = 20_000
TOTAL_INTERNAL_CHAR_LIMIT = 80_000

READABLE_EXTS = {".pdf", ".docx", ".pptx", ".txt", ".md"}

DEFAULT_STALE_MONTHS = 3

# ---------------------------------------------------------------------------
# Staleness / update checks
# ---------------------------------------------------------------------------

def get_report_path(company_dir: Path, company_name: str, output_dir: Path | None) -> Path:
    """Return the expected output path for a company's web_ingest.txt."""
    safe_name = re.sub(r"[^A-Za-z0-9._-]+", "_", company_name)[:80]
    filename = f"{safe_name}_web_ingest.txt"
    return (output_dir / filename) if output_dir else (company_dir / filename)


def report_mtime(report_path: Path) -> datetime | None:
    """Return the UTC mtime of an existing report, or None if it doesn't exist."""
    if not report_path.exists():
        return None
    return datetime.fromtimestamp(report_path.stat().st_mtime, tz=timezone.utc)


def is_stale(report_mtime_dt: datetime | None, stale_months: int) -> bool:
    """True if report doesn't exist or is older than stale_months."""
    if report_mtime_dt is None:
        return True
    age_days = (datetime.now(tz=timezone.utc) - report_mtime_dt).days
    return age_days > stale_months * 30


def has_new_internal_docs(company_dir: Path, report_mtime_dt: datetime | None) -> bool:
    """True if any readable internal doc is newer than the existing report."""
    if report_mtime_dt is None:
        return True  # no report yet — treat as new
    for path in company_dir.rglob("*"):
        if not path.is_file():
            continue
        if path.suffix.lower() not in READABLE_EXTS:
            continue
        file_mtime = datetime.fromtimestamp(path.stat().st_mtime, tz=timezone.utc)
        if file_mtime > report_mtime_dt:
            return True
    return False


def needs_update(company_dir: Path, report_path: Path, stale_months: int, force: bool) -> tuple[bool, str]:
    """
    Returns (should_run: bool, reason: str).
    Checks in order: --force > stale report > new internal docs.
    """
    if force:
        return True, "--force flag set"

    mtime = report_mtime(report_path)

    if is_stale(mtime, stale_months):
        if mtime is None:
            return True, "no existing report"
        age_days = (datetime.now(tz=timezone.utc) - mtime).days
        return True, f"report is {age_days}d old (threshold: {stale_months * 30}d)"

    if has_new_internal_docs(company_dir, mtime):
        return True, "internal docs updated since last report"

    report_age_days = (datetime.now(tz=timezone.utc) - mtime).days
    return False, f"up-to-date (report {report_age_days}d old, no new docs)"


# ---------------------------------------------------------------------------
# File extractors
# ---------------------------------------------------------------------------

def extract_pdf(path: Path) -> str:
    try:
        reader = PdfReader(str(path))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    except Exception as e:
        print(f"  [warn] PDF extract failed for {path.name}: {e}")
        return ""


def extract_docx(path: Path) -> str:
    try:
        doc = Document(str(path))
        parts = [p.text for p in doc.paragraphs if p.text]
        for tbl in doc.tables:
            for row in tbl.rows:
                cells = [c.text.strip() for c in row.cells if c.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n".join(parts)
    except Exception as e:
        print(f"  [warn] DOCX extract failed for {path.name}: {e}")
        return ""


def extract_pptx(path: Path) -> str:
    try:
        prs = Presentation(str(path))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"--- Slide {i} ---")
            for shape in slide.shapes:
                try:
                    if shape.has_text_frame and shape.text_frame.text:
                        parts.append(shape.text_frame.text)
                except Exception:
                    continue
        return "\n".join(parts)
    except Exception as e:
        print(f"  [warn] PPTX extract failed for {path.name}: {e}")
        return ""


def extract_text_file(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8", errors="ignore")
    except Exception as e:
        print(f"  [warn] text read failed for {path.name}: {e}")
        return ""


EXTRACTORS = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".txt": extract_text_file,
    ".md": extract_text_file,
}


def collect_internal_docs(company_dir: Path) -> list[dict]:
    """Returns list of {filename, text} for all readable files in the company folder."""
    docs = []
    total = 0
    for path in sorted(company_dir.rglob("*")):
        if not path.is_file():
            continue
        if len(path.relative_to(company_dir).parts) > 2:
            continue
        if path.suffix.lower() not in READABLE_EXTS:
            continue
        text = EXTRACTORS[path.suffix.lower()](path)
        if not text.strip():
            continue
        text = text[:PER_FILE_CHAR_LIMIT]
        docs.append({"filename": path.relative_to(company_dir).as_posix(), "text": text})
        total += len(text)
        if total >= TOTAL_INTERNAL_CHAR_LIMIT:
            break
    return docs


# ---------------------------------------------------------------------------
# Claude with web_search tool
# ---------------------------------------------------------------------------

SEARCH_SYSTEM_PROMPT = """You are a medtech Venture Capital research assistant. Your only job is to search the web and return raw findings.

Search for information about the given company. Run multiple searches covering:
- Company website, about page, product details
- Recent funding rounds and investors
- Clinical trial results and regulatory status
- Executive team and founders
- News articles and press releases
- Competitor comparisons
- Published research papers

Return ONLY a numbered list of findings in this exact format, one per line:
[EXT-1] <url> | <title> | <one sentence summary of what this source says>
[EXT-2] <url> | <title> | <one sentence summary>
...

No prose. No analysis. Just the numbered list."""


REPORT_SYSTEM_PROMPT = """
You are a medtech venture capital research analyst.

Write ONE complete but concise medtech VC research report using the internal documents and web findings provided.

IMPORTANT LENGTH LIMIT:
- Prioritize the most investment-relevant information.
- Use concise bullets instead of long paragraphs.
- Do not repeat the same facts across sections.
- Limit competitor analysis to the 3-5 most relevant competitors.
- For each competitor, use no more than 4 bullets.
- The final answer must include all 9 sections and REFERENCES.

CITATION FORMAT — mandatory:
- Every factual sentence must end with an inline citation tag.
- Internal sources: [INT-1], [INT-2], etc.
- Web sources: [EXT-1], [EXT-2], etc.
- Multiple sources: [INT-1][EXT-2].
- Citations go inside the sentence before the period.
- Do not include uncited factual claims.
- Do not invent missing information; write "not found in the provided sources" when unavailable.

REPORT SECTIONS:

1. Company Overview
Include company name, founding year, headquarters, stage, disease area, target users/patients, problem solved, and a 1-sentence VC thesis.

2. Technology & Product
Include product name, device/software category, mechanism of action, workflow, site of care, key technical components, and main clinical advantages vs standard of care.

3. Clinical Trial Results
Include study type, sample size, endpoints, follow-up, key outcomes, safety, limitations, and whether evidence is preclinical, clinical, or commercial. If unavailable, say so.

4. Clinical & Regulatory Status
Include development stage, FDA/CE status, likely pathway, regulatory risks, reimbursement status, and whether evidence appears sufficient for clearance/approval.

5. Funding & Financials
Include total funding, latest round, investors, grants, revenue status, business model, pricing, capital needs, and exit relevance. If unavailable, say so.

6. Leadership Team with Contact
Include CEO/founders/key executives, relevant background, and email or LinkedIn if available. If no contact is found, write "No contact information is able to be found."

7. Competing Companies
Cover only the 3-5 most relevant direct or indirect competitors. For each competitor include:
- Product/technology
- Stage/regulatory status
- Key metrics or traction
- Why the current company may be advantageous
- Where the competitor may be stronger

8. Recent News & Developments
Include major updates from the last 24 months: funding, trials, FDA/CE, partnerships, publications, launches, hires, or commercial milestones. If unavailable, say so.

9. Information Gaps
List the most important diligence gaps, especially clinical data, regulatory pathway, reimbursement, pricing, manufacturing, IP, commercial traction, and team gaps.

End with:
REFERENCES
[INT-1] filename
[INT-2] filename
[EXT-1] Title — URL
[EXT-2] Title — URL

Writing style:
- Professional VC diligence tone.
- Analytical, not promotional.
- Concise and evidence-based.
- Do not include long explanations.
- Do not include every detail from the source material; synthesize only what matters for investment diligence.
"""


def search_web(client, company_name: str) -> str:
    """Step 1: Use Claude with web_search to gather raw findings. Returns numbered EXT list."""
    messages = [{"role": "user", "content": f"Search the web for information about this medtech company: {company_name}"}]

    while True:
        resp = client.messages.create(
            model=CLAUDE_MODEL,
            max_tokens=12000,
            system=SEARCH_SYSTEM_PROMPT,
            tools=[{"type": "web_search_20250305", "name": "web_search"}],
            messages=messages,
        )

        tool_uses = [b for b in resp.content if b.type == "tool_use"]
        text_parts = [b.text for b in resp.content if b.type == "text"]

        if resp.stop_reason == "end_turn" or not tool_uses:
            return "\n".join(text_parts).strip()

        messages.append({"role": "assistant", "content": resp.content})
        tool_results = []
        for tool_use in tool_uses:
            query = tool_use.input.get("query", "")
            print(f"    web search: {query}")
            tool_results.append({
                "type": "tool_result",
                "tool_use_id": tool_use.id,
                "content": "",
            })
        messages.append({"role": "user", "content": tool_results})


def write_report(client, company_name: str, internal_docs: list, web_findings: str) -> str:
    """Step 2: Write the cited report given internal docs + web findings."""
    parts = [f"Company: {company_name}\n"]

    parts.append("=== INTERNAL DOCUMENTS ===")
    for i, doc in enumerate(internal_docs, 1):
        parts.append(f"\n[INT-{i}] File: {doc['filename']}\n{doc['text'][:6000]}")

    parts.append("\n=== WEB FINDINGS (already searched) ===")
    parts.append(web_findings)

    parts.append("\n=== TASK ===\nWrite the full research report with inline citations as instructed.")

    resp = client.messages.create(
        model=CLAUDE_MODEL,
        max_tokens=8000,
        system=REPORT_SYSTEM_PROMPT,
        messages=[{"role": "user", "content": "\n".join(parts)}],
    )
    return "".join(b.text for b in resp.content if b.type == "text").strip()


def run_research(client, company_name: str, internal_docs: list) -> tuple:
    """Two-step: search first, then write report. Returns (report_text, [])."""
    print("  step 1: searching web...")
    web_findings = search_web(client, company_name)

    print("  step 2: writing cited report...")
    report = write_report(client, company_name, internal_docs, web_findings)
    return report, []


def load_state(path: Path) -> dict:
    """Load state as a dict mapping state_key -> last_run ISO timestamp."""
    if path.exists():
        try:
            data = json.loads(path.read_text())
            # Migrate old format (plain set/list) to new dict format
            if isinstance(data, list):
                return {k: None for k in data}
            if isinstance(data, dict):
                return data
        except Exception:
            pass
    return {}


def save_state(state: dict, path: Path) -> None:
    path.write_text(json.dumps(state, indent=2, sort_keys=True))


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main() -> None:
    ap = argparse.ArgumentParser(
        description="Web research enrichment — writes cited .txt reports per company. "
                    "Re-runs automatically when reports are stale (>3 months) or internal docs change."
    )
    ap.add_argument("--event", help="Restrict to one event folder (exact folder name).")
    ap.add_argument("--company", help="Restrict to one company folder (exact folder name).")
    ap.add_argument("--force", action="store_true",
                    help="Force re-research even if the report is up-to-date. "
                         "Combine with --company or --event to target specific companies.")
    ap.add_argument("--dry-run", action="store_true", help="Print reports to stdout, don't write files.")
    ap.add_argument("--reset-state", action="store_true", help="Clear the skip cache and exit.")
    ap.add_argument("--output-dir", help="Write all .txt files here instead of inside each company folder.")
    ap.add_argument("--root", help="Override HUNNIWELL_COMPANYFILES_ROOT.")
    ap.add_argument("--stale-months", type=int, default=DEFAULT_STALE_MONTHS,
                    help=f"Re-run if report is older than this many months (default: {DEFAULT_STALE_MONTHS}).")
    args = ap.parse_args()

    root = Path(args.root or os.environ.get("HUNNIWELL_COMPANYFILES_ROOT") or DEFAULT_ROOT).expanduser()

    if args.reset_state:
        if STATE_FILE.exists():
            STATE_FILE.unlink()
            print(f"Removed {STATE_FILE}")
        else:
            print("No state file found.")
        return

    anthropic_key = os.environ.get("ANTHROPIC_API_KEY")
    if not anthropic_key:
        print("ERROR: ANTHROPIC_API_KEY not set.")
        sys.exit(1)
    if not root.exists():
        print(f"ERROR: Root directory not found: {root}")
        sys.exit(1)

    output_dir = Path(args.output_dir).resolve() if args.output_dir else None
    if output_dir:
        output_dir.mkdir(parents=True, exist_ok=True)

    client = anthropic.Anthropic(api_key=anthropic_key)
    state = load_state(STATE_FILE)

    all_events = sorted(p for p in root.iterdir() if p.is_dir())
    if args.event:
        event_folders = [p for p in all_events if p.name == args.event]
        if not event_folders:
            print(f"No event folder matching '{args.event}'.")
            sys.exit(1)
    else:
        event_folders = all_events

    total_ok = total_skip = total_fail = 0

    for event_folder in event_folders:
        company_dirs = sorted(p for p in event_folder.iterdir() if p.is_dir())

        for company_dir in company_dirs:
            company = company_dir.name.strip()
            state_key = f"{event_folder.name}/{company}"

            if args.company and company != args.company:
                continue

            report_path = get_report_path(company_dir, company, output_dir)

            # --force only applies to the targeted scope (--company / --event / all)
            force_this = args.force

            should_run, reason = needs_update(
                company_dir, report_path, args.stale_months, force=force_this
            )

            print(f"\n[{event_folder.name}] {company}")
            if not should_run:
                print(f"  skip ({reason})")
                total_skip += 1
                continue

            print(f"  update triggered: {reason}")

            # Read internal docs
            print("  reading internal docs...")
            internal_docs = collect_internal_docs(company_dir)
            print(f"  found {len(internal_docs)} internal file(s)")

            # Claude researches + writes report
            print("  researching (Claude + web search)...")
            try:
                report, _ = run_research(client, company, internal_docs)
            except Exception as e:
                print(f"  ERROR: {e}")
                total_fail += 1
                continue

            # Build output
            header = (
                f"RESEARCH REPORT: {company}\n"
                f"Event: {event_folder.name}\n"
                f"Generated: {time.strftime('%Y-%m-%d %H:%M')}\n"
                f"Internal files: {len(internal_docs)}\n"
                f"Update reason: {reason}\n"
                + "=" * 80 + "\n\n"
            )
            full_report = header + report

            if args.dry_run:
                print("\n" + full_report)
            else:
                report_path.write_text(full_report, encoding="utf-8")
                print(f"  -> wrote {report_path}")

            # Record last-run timestamp in state
            state[state_key] = datetime.now(tz=timezone.utc).isoformat()
            save_state(state, STATE_FILE)
            total_ok += 1

    print(f"\nDone. updated={total_ok}  skipped={total_skip}  failed={total_fail}")


if __name__ == "__main__":
    main()
